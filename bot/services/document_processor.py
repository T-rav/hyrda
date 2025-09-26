"""
Document Processor Service

Handles processing and preparation of documents before ingestion.
Supports various file formats and text extraction methods.
"""

import json
import logging
import re
from io import BytesIO
from typing import Any

# Document processing libraries for binary files
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word documents
from openpyxl import load_workbook  # openpyxl for Excel
from pptx import Presentation  # python-pptx for PowerPoint

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Helper service for processing documents before ingestion"""

    def process_text_file(
        self, content: str, metadata: dict[str, Any] | None = None
    ) -> list[dict[str, Any]]:
        """
        Process plain text files into chunks.

        Args:
            content: Raw text content
            metadata: Optional metadata to include

        Returns:
            List of processed document chunks
        """
        if not content.strip():
            return []

        # Simple chunking by paragraphs for text files
        paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]

        if not paragraphs:
            # Fallback: treat entire content as single chunk
            paragraphs = [content.strip()]

        chunks = []
        base_metadata = metadata or {}

        for i, paragraph in enumerate(paragraphs):
            if len(paragraph) > 50:  # Only include substantial chunks
                chunk_metadata = base_metadata.copy()
                chunk_metadata.update(
                    {
                        "chunk_id": i,
                        "content_type": "text",
                        "chunk_size": len(paragraph),
                    }
                )

                chunks.append({"content": paragraph, "metadata": chunk_metadata})

        logger.debug(f"📄 Processed text file into {len(chunks)} chunks")
        return chunks

    def process_markdown_file(
        self, content: str, metadata: dict[str, Any] | None = None
    ) -> list[dict[str, Any]]:
        """
        Process Markdown files with section-aware chunking.

        Args:
            content: Raw markdown content
            metadata: Optional metadata to include

        Returns:
            List of processed document chunks
        """
        if not content.strip():
            return []

        chunks = []
        base_metadata = metadata or {}

        # Split by headers (basic approach)
        sections = []
        current_section = []
        current_header = None

        for line in content.split("\n"):
            stripped_line = line.strip()
            if stripped_line.startswith("#"):
                # New header found
                if current_section:
                    sections.append(
                        {
                            "header": current_header,
                            "content": "\n".join(current_section),
                        }
                    )
                current_header = stripped_line
                current_section = [stripped_line]
            else:
                current_section.append(stripped_line)

        # Don't forget the last section
        if current_section:
            sections.append(
                {"header": current_header, "content": "\n".join(current_section)}
            )

        # Convert sections to chunks
        for i, section in enumerate(sections):
            if len(section["content"]) > 100:  # Only substantial sections
                chunk_metadata = base_metadata.copy()
                chunk_metadata.update(
                    {
                        "chunk_id": i,
                        "content_type": "markdown",
                        "section_header": section.get("header"),
                        "chunk_size": len(section["content"]),
                    }
                )

                chunks.append(
                    {"content": section["content"], "metadata": chunk_metadata}
                )

        logger.debug(f"📝 Processed markdown file into {len(chunks)} sections")
        return chunks

    def process_json_file(
        self, content: str, metadata: dict[str, Any] | None = None
    ) -> list[dict[str, Any]]:
        """
        Process JSON files by extracting meaningful text content.

        Args:
            content: Raw JSON content
            metadata: Optional metadata to include

        Returns:
            List of processed document chunks
        """
        try:
            data = json.loads(content)
        except json.JSONDecodeError as e:
            logger.warning(f"Failed to parse JSON content: {e}")
            return []

        chunks = []
        base_metadata = metadata or {}

        def extract_text_from_json(obj: Any, path: str = "") -> list[str]:
            """Recursively extract text values from JSON"""
            texts = []

            if isinstance(obj, dict):
                for key, value in obj.items():
                    current_path = f"{path}.{key}" if path else key
                    if isinstance(value, str) and len(value.strip()) > 20:
                        texts.append(f"{key}: {value}")
                    else:
                        texts.extend(extract_text_from_json(value, current_path))
            elif isinstance(obj, list):
                for i, item in enumerate(obj):
                    current_path = f"{path}[{i}]" if path else f"[{i}]"
                    texts.extend(extract_text_from_json(item, current_path))
            elif isinstance(obj, str) and len(obj.strip()) > 20:
                texts.append(obj)

            return texts

        text_chunks = extract_text_from_json(data)

        # Combine small chunks and create document chunks
        combined_chunks = []
        current_chunk = []
        current_size = 0
        max_chunk_size = 1000  # characters

        for text in text_chunks:
            if current_size + len(text) > max_chunk_size and current_chunk:
                # Save current chunk and start new one
                combined_chunks.append("\n".join(current_chunk))
                current_chunk = [text]
                current_size = len(text)
            else:
                current_chunk.append(text)
                current_size += len(text)

        # Don't forget the last chunk
        if current_chunk:
            combined_chunks.append("\n".join(current_chunk))

        # Convert to final format
        for i, chunk_content in enumerate(combined_chunks):
            chunk_metadata = base_metadata.copy()
            chunk_metadata.update(
                {
                    "chunk_id": i,
                    "content_type": "json",
                    "chunk_size": len(chunk_content),
                }
            )

            chunks.append({"content": chunk_content, "metadata": chunk_metadata})

        logger.debug(f"📊 Processed JSON file into {len(chunks)} chunks")
        return chunks

    def process_generic_document(
        self, content: str, file_type: str, metadata: dict[str, Any] | None = None
    ) -> list[dict[str, Any]]:
        """
        Process any document type with generic chunking strategy.

        Args:
            content: Document content
            file_type: Type/extension of the file
            metadata: Optional metadata to include

        Returns:
            List of processed document chunks
        """
        if not content.strip():
            return []

        # Determine processing strategy based on file type
        if file_type.lower() in ["md", "markdown"]:
            return self.process_markdown_file(content, metadata)
        elif file_type.lower() == "json":
            return self.process_json_file(content, metadata)
        else:
            return self.process_text_file(content, metadata)

    def validate_document_content(self, content: str) -> dict[str, Any]:
        """
        Validate document content quality and characteristics.

        Args:
            content: Document content to validate

        Returns:
            Validation results and metrics
        """
        if not content:
            return {
                "valid": False,
                "issues": ["Content is empty"],
                "metrics": {"length": 0, "lines": 0, "words": 0},
            }

        issues = []
        metrics = {
            "length": len(content),
            "lines": len(content.split("\n")),
            "words": len(content.split()),
            "paragraphs": len([p for p in content.split("\n\n") if p.strip()]),
        }

        # Check for potential issues
        if metrics["length"] < 50:
            issues.append("Content is very short")

        if metrics["words"] < 10:
            issues.append("Very few words detected")

        # Check for mostly non-text content
        non_text_chars = sum(
            1 for c in content if not c.isprintable() and c not in "\n\r\t"
        )
        if non_text_chars > metrics["length"] * 0.1:
            issues.append("High ratio of non-printable characters")

        return {"valid": len(issues) == 0, "issues": issues, "metrics": metrics}

    def chunk_large_text(
        self, text: str, chunk_size: int = 1000, overlap: int = 100
    ) -> list[str]:
        """
        Chunk large text into smaller pieces with optional overlap.

        Args:
            text: Text to chunk
            chunk_size: Maximum size of each chunk
            overlap: Number of characters to overlap between chunks

        Returns:
            List of text chunks
        """
        if not text:
            return [""]

        if len(text) <= chunk_size:
            return [text]

        chunks = []
        start = 0

        while start < len(text):
            end = start + chunk_size

            # If this isn't the last chunk, try to break at word boundary
            if end < len(text):
                # Look for word boundary within the last 100 characters
                word_boundary = text.rfind(" ", max(start, end - 100), end)
                if word_boundary > start:
                    end = word_boundary

            chunk = text[start:end]
            chunks.append(chunk)

            # Move start forward, accounting for overlap
            start = end - overlap if overlap > 0 else end

            # Prevent infinite loop
            if start >= len(text):
                break

        return chunks

    def extract_metadata_from_content(
        self, content: str, content_type: str
    ) -> dict[str, Any]:
        """
        Extract metadata from document content.

        Args:
            content: Document content
            content_type: Type of content (text, markdown, etc.)

        Returns:
            Extracted metadata
        """
        metadata = {"content_type": content_type}

        if content_type == "markdown":
            # Extract title from first header
            lines = content.split("\n")
            for original_line in lines:
                stripped_line = original_line.strip()
                if stripped_line.startswith("#"):
                    # Remove # symbols and get title
                    title = stripped_line.lstrip("#").strip()
                    if title:
                        metadata["extracted_title"] = title
                        break
        else:
            # Look for title patterns in text
            lines = content.split("\n")
            for original_line in lines:
                stripped_line = original_line.strip()
                if stripped_line.lower().startswith("title:"):
                    title = stripped_line[6:].strip()  # Remove 'title:' prefix
                    if title:
                        metadata["extracted_title"] = title
                        break

        return metadata

    def clean_text_content(self, text: str) -> str:
        """
        Clean and normalize text content.

        Args:
            text: Raw text content

        Returns:
            Cleaned text
        """
        if not text:
            return ""

        # Replace tabs and carriage returns with spaces
        text = text.replace("\t", " ").replace("\r", "")

        # Remove null characters and other control characters
        text = "".join(char for char in text if char.isprintable() or char in "\n ")

        # Normalize multiple spaces to single spaces, but preserve line breaks
        text = re.sub(r"[ ]+", " ", text)

        # Strip leading/trailing whitespace
        text = text.strip()

        return text

    def get_content_summary(self, content: str, max_length: int = 200) -> str:
        """
        Generate a summary of content.

        Args:
            content: Content to summarize
            max_length: Maximum length of summary

        Returns:
            Content summary
        """
        if not content:
            return ""

        if len(content) <= max_length:
            return content

        # Try to break at sentence boundary
        sentences = content.split(". ")
        summary = ""

        for sentence in sentences:
            if len(summary + sentence + ". ") <= max_length - 3:
                summary += sentence + ". "
            else:
                break

        if not summary:
            # Fallback to simple truncation
            summary = content[: max_length - 3]

        return summary.strip() + "..."

    def validate_document_structure(self, document: dict[str, Any]) -> bool:
        """
        Validate document structure.

        Args:
            document: Document to validate

        Returns:
            True if document is valid
        """
        if not isinstance(document, dict):
            return False

        if "content" not in document:
            return False

        content = document["content"]
        if not content or not isinstance(content, str):
            return False

        # Check minimum content length
        return len(content.strip()) >= 10

    def extract_text(self, content: bytes, mime_type: str) -> str | None:
        """
        Extract text from document bytes based on MIME type.
        Used for processing files from Slack chat.

        Args:
            content: Document content as bytes
            mime_type: MIME type of the document

        Returns:
            Extracted text content, or None if extraction fails
        """
        # Map mime types to extraction methods
        extractors = {
            "application/pdf": self._extract_pdf_text,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._extract_docx_text,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._extract_xlsx_text,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._extract_pptx_text,
        }

        try:
            # Check for direct mime type match
            if mime_type in extractors:
                return extractors[mime_type](content)

            # Check for text-based mime types
            if mime_type.startswith("text/"):
                return self._extract_text_content(content)

            # Unsupported mime type
            logger.warning(f"Unsupported MIME type for text extraction: {mime_type}")
            return None
        except Exception as e:
            logger.error(f"Error extracting text from {mime_type}: {e}")
            return None

    def _extract_text_content(self, content: bytes) -> str | None:
        """Extract text from plain text content."""
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("latin-1", errors="ignore")

    def _extract_pdf_text(self, pdf_content: bytes) -> str | None:
        """
        Extract text content from PDF bytes using PyMuPDF.

        Args:
            pdf_content: PDF file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PDF from bytes
            pdf_stream = BytesIO(pdf_content)
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

            # Extract text from all pages
            text_content = []
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append(text)

            pdf_document.close()

            # Join all pages with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return None

    def _extract_docx_text(self, docx_content: bytes) -> str | None:
        """
        Extract text content from Word document bytes using python-docx.

        Args:
            docx_content: Word document file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Word document from bytes
            docx_stream = BytesIO(docx_content)
            doc = Document(docx_stream)

            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:  # Only add non-empty paragraphs
                    text_content.append(text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_content.append(" | ".join(row_text))

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            logger.error(f"Error extracting Word document text: {e}")
            return None

    def _extract_xlsx_text(self, xlsx_content: bytes) -> str | None:
        """
        Extract text content from Excel spreadsheet bytes using openpyxl.

        Args:
            xlsx_content: Excel spreadsheet file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Excel workbook from bytes
            xlsx_stream = BytesIO(xlsx_content)
            workbook = load_workbook(xlsx_stream, read_only=True, data_only=True)

            text_content = []

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]

                # Extract text from all cells with data
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell_value in row:
                        if cell_value is not None:
                            # Convert to string and clean up
                            cell_text = str(cell_value).strip()
                            if cell_text:
                                row_data.append(cell_text)

                    if row_data:  # Only add non-empty rows
                        sheet_content.append(" | ".join(row_data))

                if len(sheet_content) > 1:  # More than just the sheet name
                    text_content.extend(sheet_content)

            workbook.close()

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            logger.error(f"Error extracting Excel spreadsheet text: {e}")
            return None

    def _extract_pptx_text(self, pptx_content: bytes) -> str | None:
        """
        Extract text content from PowerPoint presentation bytes using python-pptx.

        Args:
            pptx_content: PowerPoint presentation file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PowerPoint presentation from bytes
            pptx_stream = BytesIO(pptx_content)
            presentation = Presentation(pptx_stream)

            text_content = []

            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]

                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_content.append(shape_text)

                    # Handle tables in slides
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                slide_content.append(" | ".join(row_text))

                if len(slide_content) > 1:  # More than just the slide number
                    text_content.extend(slide_content)

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            logger.error(f"Error extracting PowerPoint presentation text: {e}")
            return None
