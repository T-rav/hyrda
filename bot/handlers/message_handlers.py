import contextlib
import io
import logging
import time

import requests

from handlers.constants import (
    CHUNK_OVERLAP_CHARS,
    FILE_DOWNLOAD_TIMEOUT_SECONDS,
    MAX_EMBEDDING_CHARS,
    MAX_FILE_SIZE_BYTES,
)

try:
    import fitz  # PyMuPDF  # type: ignore[reportMissingImports]

    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

try:
    from docx import Document  # type: ignore[reportMissingImports]

    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

try:
    from openpyxl import load_workbook  # type: ignore[reportMissingModuleSource]

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

try:
    from pptx import Presentation  # type: ignore[reportMissingImports]

    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

from services.formatting import MessageFormatter
from services.langfuse_service import get_langfuse_service
from services.llm_service import LLMService
from services.metrics_service import get_metrics_service
from services.prompt_service import get_prompt_service
from services.slack_service import SlackService
from utils.errors import handle_error

logger = logging.getLogger(__name__)

# Constants
HTTP_OK = 200


async def extract_pdf_text(pdf_content: bytes, file_name: str) -> str:
    """Extract text from PDF using PyMuPDF and chunk for embedding compatibility"""
    if not PYMUPDF_AVAILABLE:
        logger.error("PyMuPDF not available for PDF processing")
        return f"[PDF file: {file_name} - PyMuPDF library not available]"

    try:
        from services.embedding import chunk_text

        # Open PDF from bytes
        pdf_document = fitz.open(stream=pdf_content, filetype="pdf")
        text_content = ""

        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            page_text = page.get_text()
            if page_text.strip():
                text_content += f"\n\n--- Page {page_num + 1} ---\n{page_text}"

        pdf_document.close()

        if text_content.strip():
            full_text = text_content.strip()

            # If content is too long, chunk it to prevent embedding failures
            # Conservative limit: MAX_EMBEDDING_CHARS ≈ 1500 tokens (well under 8192 limit)
            if len(full_text) > MAX_EMBEDDING_CHARS:
                logger.info(
                    f"PDF content is {len(full_text)} chars, chunking for embedding compatibility"
                )
                chunks = chunk_text(
                    full_text,
                    chunk_size=MAX_EMBEDDING_CHARS,
                    chunk_overlap=CHUNK_OVERLAP_CHARS,
                )
                # Return first chunk with indicator
                chunked_content = chunks[0]
                if len(chunks) > 1:
                    chunked_content += (
                        f"\n\n[Note: This is chunk 1 of {len(chunks)} from {file_name}]"
                    )
                return chunked_content
            else:
                return full_text
        else:
            return f"[PDF file: {file_name} - No extractable text content found]"

    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_name}: {e}")
        return f"[PDF file: {file_name} - Error extracting text: {str(e)}]"


async def extract_office_text(content: bytes, file_name: str, file_type: str) -> str:
    """Extract text from Office documents (Word, Excel, PowerPoint)"""
    try:
        content_stream = io.BytesIO(content)

        if file_name.endswith(".docx") or "wordprocessingml" in file_type:
            return await extract_word_text(content_stream, file_name)
        elif file_name.endswith(".xlsx") or "spreadsheetml" in file_type:
            return await extract_excel_text(content_stream, file_name)
        elif file_name.endswith(".pptx") or "presentationml" in file_type:
            return await extract_powerpoint_text(content_stream, file_name)
        else:
            return f"[Office document: {file_name} - Unsupported Office format]"

    except Exception as e:
        logger.error(f"Error extracting text from Office document {file_name}: {e}")
        return f"[Office document: {file_name} - Error extracting text: {str(e)}]"


async def extract_word_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from Word documents"""
    if not PYTHON_DOCX_AVAILABLE:
        logger.error("python-docx not available for Word processing")
        return f"[Word document: {file_name} - python-docx library not available]"

    try:
        doc = Document(content_stream)
        text_content = ""

        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content += paragraph.text + "\n"

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content += " | ".join(row_text) + "\n"

        if text_content.strip():
            return text_content.strip()
        else:
            return f"[Word document: {file_name} - No extractable text content found]"

    except Exception as e:
        logger.error(f"Error extracting text from Word document {file_name}: {e}")
        return f"[Word document: {file_name} - Error extracting text: {str(e)}]"


async def extract_excel_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from Excel files"""
    if not OPENPYXL_AVAILABLE:
        logger.error("openpyxl not available for Excel processing")
        return f"[Excel file: {file_name} - openpyxl library not available]"

    try:
        workbook = load_workbook(content_stream, read_only=True, data_only=True)
        text_content = ""

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"\n\n--- Sheet: {sheet_name} ---\n"

            # Extract data from used cells
            for row in sheet.iter_rows(values_only=True):
                row_data = []
                for cell_value in row:
                    if cell_value is not None:
                        row_data.append(str(cell_value))
                if row_data:
                    text_content += " | ".join(row_data) + "\n"

        workbook.close()

        if text_content.strip():
            return text_content.strip()
        else:
            return f"[Excel file: {file_name} - No extractable data found]"

    except Exception as e:
        logger.error(f"Error extracting text from Excel file {file_name}: {e}")
        return f"[Excel file: {file_name} - Error extracting text: {str(e)}]"


async def extract_powerpoint_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from PowerPoint presentations"""
    if not PYTHON_PPTX_AVAILABLE:
        logger.error("python-pptx not available for PowerPoint processing")
        return f"[PowerPoint file: {file_name} - python-pptx library not available]"

    try:
        prs = Presentation(content_stream)
        text_content = ""

        for i, slide in enumerate(prs.slides):
            text_content += f"\n\n--- Slide {i + 1} ---\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content += shape.text + "\n"

        if text_content.strip():
            return text_content.strip()
        else:
            return f"[PowerPoint file: {file_name} - No extractable text content found]"

    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint file {file_name}: {e}")
        return f"[PowerPoint file: {file_name} - Error extracting text: {str(e)}]"


async def process_file_attachments(
    files: list[dict], slack_service: SlackService
) -> str:
    """Process file attachments and extract text content"""
    document_content = ""

    for file_info in files:
        try:
            file_name = file_info.get("name", "unknown")
            file_type = file_info.get("mimetype", "")
            file_size = file_info.get("size", 0)

            # Skip very large files (>100MB)
            if file_size > MAX_FILE_SIZE_BYTES:
                logger.warning(f"Skipping large file: {file_name} ({file_size} bytes)")
                continue

            # Get file URL
            file_url = file_info.get("url_private")
            if not file_url:
                logger.warning(f"No private URL for file: {file_name}")
                continue

            logger.info(
                f"Processing file: {file_name} ({file_type}, {file_size} bytes)"
            )

            # Download file content
            # Use 10-minute timeout for large files (up to 100MB allowed)
            headers = {"Authorization": f"Bearer {slack_service.settings.bot_token}"}
            response = requests.get(
                file_url, headers=headers, timeout=FILE_DOWNLOAD_TIMEOUT_SECONDS
            )

            if response.status_code != 200:
                logger.error(
                    f"Failed to download file {file_name}: {response.status_code}"
                )
                continue

            file_content = ""

            # Process based on file type
            # Check MIME types first, then fall back to extensions
            if (
                file_type.startswith("text/")
                or file_type in ("text/vtt", "application/x-subrip", "text/srt")
                or file_name.endswith(
                    (".txt", ".md", ".py", ".js", ".json", ".csv", ".vtt", ".srt")
                )
            ):
                # Text files
                try:
                    file_content = response.text
                except UnicodeDecodeError:
                    file_content = response.content.decode("utf-8", errors="ignore")

            elif file_type == "application/pdf" or file_name.endswith(".pdf"):
                # PDF text extraction using PyMuPDF
                file_content = await extract_pdf_text(response.content, file_name)

            elif file_type.startswith(
                "application/vnd.openxmlformats"
            ) or file_name.endswith((".docx", ".xlsx", ".pptx")):
                # Office documents
                file_content = await extract_office_text(
                    response.content, file_name, file_type
                )

            else:
                # Unknown file type
                file_content = f"[File: {file_name} ({file_type}) - content extraction not supported for this file type.]"

            if file_content:
                document_content += f"\n\n--- Content from {file_name} ---\n{file_content}\n--- End of {file_name} ---"

        except Exception as e:
            logger.error(
                f"Error processing file {file_info.get('name', 'unknown')}: {e}"
            )
            continue

    return document_content.strip()


def get_user_system_prompt(user_id: str | None = None) -> str:
    """
    Get the system prompt with user context injected.

    Args:
        user_id: Slack user ID to look up and inject context for

    Returns:
        System prompt with user context

    """
    prompt_service = get_prompt_service()
    base_prompt = ""

    if prompt_service:
        base_prompt = prompt_service.get_system_prompt()
    else:
        # Ultimate fallback if PromptService is not available
        logger.warning("PromptService not available, using minimal fallback prompt")
        base_prompt = "I'm Insight Mesh, your AI assistant. I help you find information from your organization's knowledge base and provide intelligent assistance with your questions."

    # Inject user context if user_id provided
    if user_id:
        try:
            from services.user_service import get_user_service

            user_service = get_user_service()
            user_info = user_service.get_user_info(user_id)

            if user_info:
                user_context = "\n\n**Current User Context:**\n"
                user_context += f"- Name: {user_info.get('real_name') or user_info.get('display_name', 'Unknown')}\n"
                if user_info.get("email_address"):
                    user_context += f"- Email: {user_info['email_address']}\n"

                user_context += "\nWhen responding, you can personalize your responses knowing who the user is. Address them by name when appropriate."

                return base_prompt + user_context
        except Exception as e:
            logger.error(f"Error injecting user context: {e}")

    return base_prompt


# Thread tracking service (Redis-backed with in-memory fallback)
from services.thread_tracking import get_thread_tracking

_thread_tracking = get_thread_tracking()


# ============================================================================
# Helper functions for handle_bot_command (refactored for maintainability)
# ============================================================================


async def _handle_exit_command(
    text: str, thread_ts: str | None, slack_service: SlackService, channel: str
) -> bool:
    """Handle exit commands to clear thread tracking.

    Args:
        text: User message text
        thread_ts: Thread timestamp
        slack_service: Slack service for sending messages
        channel: Channel ID

    Returns:
        True if exit command was handled, False otherwise
    """
    exit_commands = ["exit", "stop", "done", "end", "clear"]
    if not thread_ts or not any(text.strip().lower() == cmd for cmd in exit_commands):
        return False

    if await _thread_tracking.clear_thread(thread_ts):
        try:
            await slack_service.send_message(
                channel=channel,
                text="✅ Exited agent mode. I'm back to general mode!",
                thread_ts=thread_ts,
            )
        except Exception as e:
            logger.warning(f"Failed to send exit confirmation: {e}")
        return True

    return False


async def _get_thread_agent(
    thread_ts: str | None, check_thread_context: bool
) -> str | None:
    """Get agent name associated with thread if context checking is enabled.

    Args:
        thread_ts: Thread timestamp
        check_thread_context: Whether to check thread context

    Returns:
        Agent name if thread is tracked, None otherwise
    """
    if not check_thread_context or not thread_ts:
        return None

    agent_name = await _thread_tracking.get_thread_agent(thread_ts)
    if agent_name:
        logger.info(
            f"🔗 Thread {thread_ts} belongs to agent '{agent_name}' - routing automatically"
        )
    return agent_name


def _clean_markdown_from_text(text: str) -> str:
    """Remove Slack markdown formatting from text.

    Slack sends *bold* and _italic_ which breaks command parsing.

    Args:
        text: Raw text with potential markdown

    Returns:
        Cleaned text without markdown characters
    """
    clean_text = text.strip()
    # Remove leading/trailing markdown characters
    while clean_text and clean_text[0] in ["*", "_", "~"]:
        clean_text = clean_text[1:]
    while clean_text and clean_text[-1] in ["*", "_", "~"]:
        clean_text = clean_text[:-1]
    return clean_text


async def _check_and_notify_unavailable_agent(
    text: str, slack_service: SlackService, channel: str, thread_ts: str | None
) -> bool:
    """Check if a disabled/unavailable agent was requested and notify user.

    Args:
        text: Cleaned command text
        slack_service: Slack service for sending messages
        channel: Channel ID
        thread_ts: Thread timestamp

    Returns:
        True if unavailable agent was found and user notified, False otherwise
    """
    import re

    from services.agent_registry import check_agent_availability

    # Try to parse command name from text
    match = re.match(r"^[@-]?(\w+)[\s:].*", text.strip(), re.IGNORECASE)
    if not match:
        return False

    attempted_command = match.group(1).lower()
    logger.info(f"Extracted attempted command: '{attempted_command}'")

    # Check if agent exists but is unavailable
    availability = check_agent_availability(attempted_command)
    if not availability:
        return False

    # Agent exists but is not available in Slack
    error_msg = (
        f"❌ Agent `{attempted_command}` is not available.\n\n{availability['reason']}"
    )
    logger.info(
        f"Agent '{attempted_command}' exists but unavailable: {availability['reason']}"
    )

    try:
        await slack_service.send_message(
            channel=channel,
            text=error_msg,
            thread_ts=thread_ts,
        )
    except Exception as e:
        logger.error(f"Error sending unavailable agent message: {e}")

    return True


async def _execute_agent_with_streaming(
    primary_name: str,
    query: str,
    context: dict,
    slack_service: SlackService,
    channel: str,
    thread_ts: str | None,
) -> tuple[str, dict, str | None]:
    """Execute agent via HTTP with streaming status updates.

    Args:
        primary_name: Primary agent name
        query: User query
        context: Context dict with thread_id, user_id, etc.
        slack_service: Slack service for status updates
        channel: Channel ID
        thread_ts: Thread timestamp

    Returns:
        Tuple of (response, metadata, thinking_message_ts)
    """
    from services.agent_client import get_agent_client

    agent_client = get_agent_client()

    # Send thinking indicator
    thinking_message_ts = None
    try:
        thinking_message_ts = await slack_service.send_thinking_indicator(
            channel, thread_ts
        )
    except Exception as e:
        logger.error(f"Error posting thinking message: {e}")

    logger.info(
        f"Streaming agent-service for '{primary_name}' with thread_id={context.get('thread_id')}"
    )

    response = ""
    metadata = {}

    async for event in agent_client.stream(primary_name, query, context):
        phase = event.get("phase")
        step = event.get("step")

        # Update status message in place
        if phase and thinking_message_ts and phase == "running" and step:
            status_text = f"⏳ *Running:* {step}"
            try:
                await slack_service.update_message(
                    channel, thinking_message_ts, status_text
                )
            except Exception as e:
                logger.warning(f"Failed to update status message: {e}")

        # Final response with full agent output
        if "response" in event:
            response = event.get("response", "")
            metadata = event.get("metadata", {})

    return response, metadata, thinking_message_ts


async def _track_thread_for_continuity(
    thread_ts: str | None, primary_name: str | None, metadata: dict
) -> None:
    """Track thread for future continuity or clear if agent requests it.

    Args:
        thread_ts: Thread timestamp
        primary_name: Primary agent name
        metadata: Agent response metadata
    """
    if not thread_ts or not primary_name:
        return

    # Check if agent wants to auto-clear after completion
    if metadata.get("clear_thread_tracking", False):
        await _thread_tracking.clear_thread(thread_ts)
        logger.info(
            f"🔓 Cleared both agent tracking AND LangGraph checkpoint for thread {thread_ts}"
        )
    else:
        await _thread_tracking.track_thread(thread_ts, primary_name)
        logger.info(
            f"📌 Thread {thread_ts} tracked for agent '{primary_name}' (both Redis + LangGraph checkpoint)"
        )


async def _send_agent_response(
    response: str,
    thinking_message_ts: str | None,
    slack_service: SlackService,
    channel: str,
    thread_ts: str | None,
    user_id: str,
    query: str,
    conversation_id: str | None,
) -> None:
    """Format and send agent response, track in Langfuse.

    Args:
        response: Agent response text
        thinking_message_ts: Thinking message to delete
        slack_service: Slack service
        channel: Channel ID
        thread_ts: Thread timestamp
        user_id: User ID
        query: Original query
        conversation_id: Conversation ID for tracking
    """
    # Clean up thinking message
    if thinking_message_ts:
        try:
            await slack_service.delete_thinking_indicator(channel, thinking_message_ts)
        except Exception as e:
            logger.warning(f"Error deleting thinking message: {e}")

    # Format and send response (only if non-empty)
    if response:
        formatted_response = await MessageFormatter.format_message(response)
        await slack_service.send_message(
            channel=channel, text=formatted_response, thread_ts=thread_ts
        )
    else:
        logger.info("Agent returned empty response (likely already posted message)")

    # Record agent conversation turn in Langfuse for lifetime stats
    langfuse_service = get_langfuse_service()
    if langfuse_service and response:
        try:
            langfuse_service.trace_conversation(
                user_id=user_id,
                conversation_id=conversation_id,
                user_message=query,
                bot_response=response,
            )
        except Exception as e:
            logger.warning(f"Error tracing agent conversation turn to Langfuse: {e}")


# ============================================================================
# Helper functions for handle_message (refactored for maintainability)
# ============================================================================


def _determine_conversation_id(
    channel: str, thread_ts: str | None, message_ts: str | None
) -> str:
    """Determine unique conversation ID from Slack context.

    Args:
        channel: Channel ID
        thread_ts: Thread timestamp
        message_ts: Message timestamp

    Returns:
        Unique conversation ID

    Raises:
        ValueError: If conversation ID cannot be determined
    """
    is_dm = channel.startswith("D")

    if thread_ts:
        return thread_ts  # Threaded conversation
    elif is_dm:
        return channel  # DM conversation (channel is unique per user)
    elif message_ts:
        return message_ts  # Non-threaded mention (each is separate)
    else:
        raise ValueError("Must provide thread_ts or message_ts for non-DM messages")


def _record_message_metrics(
    metrics_service,
    user_id: str,
    channel: str,
    text: str,
    conversation_id: str,
    thread_ts: str | None,
) -> None:
    """Record message processing metrics.

    Args:
        metrics_service: Metrics service instance
        user_id: User ID
        channel: Channel ID
        text: Message text
        conversation_id: Conversation ID
        thread_ts: Thread timestamp
    """
    if not metrics_service:
        return

    metrics_service.record_message_processed(
        user_id=user_id, channel_type="dm" if channel.startswith("D") else "channel"
    )

    metrics_service.record_conversation_activity(conversation_id)

    # Categorize query type for metrics
    query_category = "question"  # Default
    if text.lower().startswith(("run ", "execute ", "start ")):
        query_category = "command"
    elif thread_ts:  # In a thread, likely continuing conversation
        query_category = "conversation"

    has_context = bool(thread_ts)
    metrics_service.record_query_type(query_category, has_context)


async def _get_or_store_document_content(
    conversation_cache,
    thread_ts: str | None,
    document_content: str,
    document_filename: str | None,
) -> tuple[str, str | None]:
    """Get cached document or store new one.

    Args:
        conversation_cache: Conversation cache instance
        thread_ts: Thread timestamp
        document_content: Current document content
        document_filename: Current document filename

    Returns:
        Tuple of (document_content, document_filename)
    """
    if not thread_ts or not conversation_cache:
        return document_content, document_filename

    # Store new document if provided
    if document_content:
        await conversation_cache.store_document_content(
            thread_ts, document_content, document_filename or "unknown"
        )
        return document_content, document_filename

    # Retrieve cached document if no new document
    stored_content, stored_filename = await conversation_cache.get_document_content(
        thread_ts
    )
    if stored_content:
        logger.info(
            f"Retrieved previously stored document for thread {thread_ts}: {stored_filename}"
        )
        return stored_content, stored_filename

    return document_content, document_filename


async def _add_document_reference_to_cache(
    conversation_cache,
    thread_ts: str | None,
    document_content: str,
    files: list[dict] | None,
) -> None:
    """Add document reference summary to conversation cache.

    Args:
        conversation_cache: Conversation cache instance
        thread_ts: Thread timestamp
        document_content: Document content
        files: File list
    """
    if not (document_content and conversation_cache and thread_ts):
        return

    # Create concise summary for chat history
    file_names = [f.get("name", "unknown") for f in files] if files else ["document"]
    file_summary = f"[User shared {len(file_names)} file(s): {', '.join(file_names)} - content processed and analyzed]"

    document_message = {
        "role": "user",
        "content": file_summary,
    }
    await conversation_cache.update_conversation(
        thread_ts, document_message, is_bot_message=False
    )
    logger.info(
        f"Added document reference to conversation cache for thread {thread_ts}: {file_names}"
    )


async def _determine_rag_usage(
    conversation_cache, thread_ts: str | None
) -> tuple[bool, str | None]:
    """Determine whether to use RAG based on thread type.

    Profile threads disable RAG to avoid retrieving irrelevant documents.

    Args:
        conversation_cache: Conversation cache instance
        thread_ts: Thread timestamp

    Returns:
        Tuple of (use_rag, thread_type)
    """
    thread_type = None
    if thread_ts and conversation_cache:
        thread_type = await conversation_cache.get_thread_type(thread_ts)
        logger.info(
            f"Thread {thread_ts}: type={thread_type}, has_cache={conversation_cache is not None}"
        )

    is_profile_thread = thread_type == "profile"
    use_rag = not is_profile_thread

    if is_profile_thread:
        logger.info(f"✅ Disabling RAG for profile thread {thread_ts}")
    else:
        logger.info(
            f"📚 RAG enabled for thread {thread_ts} (thread_type={thread_type})"
        )

    return use_rag, thread_type


async def _send_llm_response(
    response: str,
    thinking_message_ts: str | None,
    slack_service: SlackService,
    channel: str,
    thread_ts: str | None,
    user_id: str,
    text: str,
    conversation_id: str,
) -> None:
    """Clean up thinking indicator and send formatted response.

    Args:
        response: LLM response text
        thinking_message_ts: Thinking message to delete
        slack_service: Slack service
        channel: Channel ID
        thread_ts: Thread timestamp
        user_id: User ID
        text: Original user message
        conversation_id: Conversation ID
    """
    # Clean up thinking message
    if thinking_message_ts:
        try:
            await slack_service.delete_thinking_indicator(channel, thinking_message_ts)
        except Exception as e:
            logger.warning(f"Error deleting thinking message: {e}")

    # Format and send response
    if response:
        formatted_response = await MessageFormatter.format_message(response)
        await slack_service.send_message(
            channel=channel, text=formatted_response, thread_ts=thread_ts
        )
    else:
        response = "I apologize, but I couldn't generate a response. Please try again."
        await slack_service.send_message(
            channel=channel,
            text=response,
            thread_ts=thread_ts,
        )

    # Record conversation turn in Langfuse
    langfuse_service = get_langfuse_service()
    if langfuse_service:
        try:
            langfuse_service.trace_conversation(
                user_id=user_id,
                conversation_id=conversation_id,
                user_message=text,
                bot_response=response or "",
            )
        except Exception as e:
            logger.warning(f"Error tracing conversation turn to Langfuse: {e}")


# ============================================================================
# Main command handler (now much cleaner!)
# ============================================================================


async def handle_bot_command(
    text: str,
    user_id: str,
    slack_service: SlackService,
    channel: str,
    thread_ts: str | None = None,
    files: list[dict] | None = None,
    document_content: str | None = None,
    llm_service: LLMService | None = None,
    check_thread_context: bool = False,
    conversation_cache=None,
    conversation_id: str | None = None,
) -> bool:
    """Handle bot agent commands using router pattern.

    Routes commands like /profile and /meddic to registered agents.
    Supports thread continuity - once an agent starts a thread, subsequent
    messages in that thread automatically route to the same agent.

    Args:
        text: Full message text (e.g., "/profile tell me about Charlotte")
        user_id: Slack user ID
        slack_service: Slack service for sending messages
        channel: Channel ID
        thread_ts: Thread timestamp if in a thread
        files: Optional list of file attachments
        document_content: Optional processed file content
        llm_service: Optional LLM service for agent use
        check_thread_context: If True, check if thread belongs to an agent
        conversation_cache: Cache for conversation documents
        conversation_id: Conversation ID for tracking

    Returns:
        True if bot command was handled, False otherwise
    """
    logger.info(
        f"handle_bot_command called with text: '{text}', check_thread_context: {check_thread_context}"
    )

    # Handle exit commands first
    if await _handle_exit_command(text, thread_ts, slack_service, channel):
        return True

    # Check if thread already belongs to an agent
    agent_name = await _get_thread_agent(thread_ts, check_thread_context)
    if agent_name:
        from services.agent_registry import get_agent_info

        agent_info = get_agent_info(agent_name)
        if agent_info:
            primary_name = agent_name
            query = text  # Use full text as query (no command prefix needed)
        else:
            logger.warning(
                f"Agent '{agent_name}' not found in registry, falling back to normal routing"
            )
            agent_info = None
            primary_name = None
            query = ""
    else:
        # Clean markdown and route command
        clean_text = _clean_markdown_from_text(text)
        logger.info(f"Cleaned text for routing: '{clean_text}'")

        from services.agent_registry import route_command

        agent_info, query, primary_name = route_command(clean_text)

    logger.info(
        f"Router results: agent_info={agent_info is not None}, primary_name={primary_name}, query='{query}'"
    )

    # Check if agent is unavailable
    if not agent_info or not primary_name:
        logger.info("No agent found in enabled registry, checking availability...")
        clean_text = _clean_markdown_from_text(text)
        if await _check_and_notify_unavailable_agent(
            clean_text, slack_service, channel, thread_ts
        ):
            return True

        logger.info("No agent found, returning False")
        return False

    logger.info(f"Routing to agent '{primary_name}' with query: {query}")

    # Build context for agent
    thread_id = f"slack_{channel}_{thread_ts}" if thread_ts else f"slack_{channel}_dm"
    context = {
        "thread_id": thread_id,  # CRITICAL: For SQLite checkpointing
        "user_id": user_id,
        "channel": channel,
        "thread_ts": thread_ts,
    }

    # Add file information if available
    if document_content:
        context["document_content"] = document_content
    if files:
        context["files"] = files
        context["file_names"] = [f.get("name", "unknown") for f in files]

    # Execute agent with streaming
    try:
        response, metadata, thinking_message_ts = await _execute_agent_with_streaming(
            primary_name, query, context, slack_service, channel, thread_ts
        )

        # Track thread for future continuity
        await _track_thread_for_continuity(thread_ts, primary_name, metadata)

        # Send response
        await _send_agent_response(
            response,
            thinking_message_ts,
            slack_service,
            channel,
            thread_ts,
            user_id,
            query,
            conversation_id,
        )

        return True

    except Exception as e:
        logger.error(f"Bot command '{primary_name}' failed: {e}")
        error_response = f"❌ Bot command '/{primary_name}' failed: {str(e)}"
        await slack_service.send_message(
            channel=channel, text=error_response, thread_ts=thread_ts
        )
        return True


async def handle_message(
    text: str,
    user_id: str,
    slack_service: SlackService,
    llm_service: LLMService,
    channel: str,
    thread_ts: str | None = None,
    files: list[dict] | None = None,
    conversation_cache=None,
    message_ts: str | None = None,
):
    """Handle an incoming message from Slack.

    Args:
        text: Message text
        user_id: Slack user ID
        slack_service: Slack service instance
        llm_service: LLM service instance
        channel: Channel ID
        thread_ts: Thread timestamp (optional)
        files: Attached files (optional)
        conversation_cache: Conversation cache instance (optional)
        message_ts: Message timestamp - used as unique conversation ID when not in thread
    """
    start_time = time.time()

    # Determine unique conversation ID
    conversation_id = _determine_conversation_id(channel, thread_ts, message_ts)

    logger.info(
        "Processing user message",
        extra={
            "user_message": text,
            "user_id": user_id,
            "channel_id": channel,
            "thread_ts": thread_ts,
            "message_ts": message_ts,
            "conversation_id": conversation_id,
            "event_type": "user_message",
        },
    )

    # Record message processing metrics
    metrics_service = get_metrics_service()
    _record_message_metrics(
        metrics_service, user_id, channel, text, conversation_id, thread_ts
    )

    # For tracking the thinking indicator message
    thinking_message_ts = None

    try:
        # Process file attachments first (needed for both bot commands and LLM)
        document_content = ""
        document_filename = None

        if files:
            logger.info(f"Files attached: {[f.get('name', 'unknown') for f in files]}")
            document_content = await process_file_attachments(files, slack_service)
            document_filename = files[0].get("name") if files else None

            if document_content:
                logger.info(
                    f"Extracted {len(document_content)} characters from {len(files)} file(s)"
                )

        # Check for bot agent commands: -profile, profile, -meddic, meddic, etc.
        # Router handles parsing and validation internally
        # Pass check_thread_context=True to enable thread continuity
        handled = await handle_bot_command(
            text=text,
            user_id=user_id,
            slack_service=slack_service,
            channel=channel,
            thread_ts=thread_ts,
            files=files,
            document_content=document_content,
            llm_service=llm_service,
            check_thread_context=True,  # Enable thread-aware routing
            conversation_cache=conversation_cache,  # Pass cache for agent-generated docs
            conversation_id=conversation_id,  # Pass for Langfuse tracking
        )

        if handled:
            return

        # Regular LLM response handling
        try:
            thinking_message_ts = await slack_service.send_thinking_indicator(
                channel, thread_ts
            )
        except Exception as e:
            logger.error(f"Error posting thinking message: {e}")

        # Store or retrieve document content from cache
        document_content, document_filename = await _get_or_store_document_content(
            conversation_cache, thread_ts, document_content, document_filename
        )

        # Get the thread history for context
        history, should_use_thread_context = await slack_service.get_thread_history(
            channel, thread_ts
        )

        if should_use_thread_context:
            logger.info(
                f"Using thread context with {len(history)} messages for response generation"
            )

        # Add document reference to cache
        await _add_document_reference_to_cache(
            conversation_cache, thread_ts, document_content, files
        )

        # Determine whether to use RAG based on thread type
        use_rag, thread_type = await _determine_rag_usage(conversation_cache, thread_ts)

        logger.info(
            f"Thread {thread_ts}: thread_type={thread_type}, use_rag={use_rag}, "
            f"has_document={document_content is not None}"
        )

        # Generate response using LLM service with document-aware RAG
        response = await llm_service.get_response(
            messages=history,
            user_id=user_id,
            current_query=text,
            document_content=document_content if document_content else None,
            document_filename=document_filename,
            conversation_id=conversation_id,
            conversation_cache=conversation_cache,
            use_rag=use_rag,
        )

        # Send formatted response and clean up thinking indicator
        await _send_llm_response(
            response,
            thinking_message_ts,
            slack_service,
            channel,
            thread_ts,
            user_id,
            text,
            conversation_id,
        )

        # Record successful request timing
        if metrics_service:
            duration = time.time() - start_time
            metrics_service.request_duration.labels(
                endpoint="message_handler", method="POST"
            ).observe(duration)

    except Exception as e:
        # Clean up thinking message on error
        if thinking_message_ts:
            with contextlib.suppress(Exception):
                await slack_service.delete_thinking_indicator(
                    channel, thinking_message_ts
                )

        await handle_error(
            slack_service.client,
            channel,
            thread_ts,
            e,
            "I'm sorry, I encountered an error while processing your message.",
        )
