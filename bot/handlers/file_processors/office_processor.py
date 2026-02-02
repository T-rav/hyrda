"""Microsoft Office document processing (Word, Excel, PowerPoint)."""

import io
import logging

try:
    from docx import Document  # type: ignore[reportMissingImports]

    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

try:
    from openpyxl import load_workbook  # type: ignore[reportMissingModuleSource]

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

try:
    from pptx import Presentation  # type: ignore[reportMissingImports]

    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False  # type: ignore[reportConstantRedefinition]

logger = logging.getLogger(__name__)


async def extract_word_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from Word document."""
    if not PYTHON_DOCX_AVAILABLE:
        logger.warning("python-docx not available")
        return f"[Word file: {file_name} - python-docx not installed]"

    try:
        doc = Document(content_stream)
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)

        if text_parts:
            return "\n\n".join(text_parts)
        else:
            return f"[Word file: {file_name} - No extractable text found]"

    except Exception as e:
        logger.error(f"Error extracting text from Word document {file_name}: {e}")
        return f"[Word file: {file_name} - Error: {str(e)}]"


async def extract_excel_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from Excel spreadsheet."""
    if not OPENPYXL_AVAILABLE:
        logger.warning("openpyxl not available")
        return f"[Excel file: {file_name} - openpyxl not installed]"

    try:
        workbook = load_workbook(content_stream, read_only=True, data_only=True)
        text_parts = []

        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(row_values):
                    text_parts.append(" | ".join(row_values))

        if text_parts:
            return "\n".join(text_parts)
        else:
            return f"[Excel file: {file_name} - No extractable data found]"

    except Exception as e:
        logger.error(f"Error extracting text from Excel file {file_name}: {e}")
        return f"[Excel file: {file_name} - Error: {str(e)}]"


async def extract_powerpoint_text(content_stream: io.BytesIO, file_name: str) -> str:
    """Extract text from PowerPoint presentation."""
    if not PYTHON_PPTX_AVAILABLE:
        logger.warning("python-pptx not available")
        return f"[PowerPoint file: {file_name} - python-pptx not installed]"

    try:
        presentation = Presentation(content_stream)
        text_parts = []

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():  # type: ignore[attr-defined]
                    slide_text.append(shape.text)  # type: ignore[attr-defined]

            if slide_text:
                text_parts.append(f"\n=== Slide {slide_num} ===\n")
                text_parts.append("\n".join(slide_text))

        if text_parts:
            return "\n".join(text_parts)
        else:
            return f"[PowerPoint file: {file_name} - No extractable text found]"

    except Exception as e:
        logger.error(f"Error extracting text from PowerPoint file {file_name}: {e}")
        return f"[PowerPoint file: {file_name} - Error: {str(e)}]"


async def extract_office_text(content: bytes, file_name: str, file_type: str) -> str:
    """Extract text from Microsoft Office files (.docx, .xlsx, .pptx)."""
    content_stream = io.BytesIO(content)

    file_type_lower = file_type.lower()
    if "wordprocessingml" in file_type_lower or file_type_lower == "docx":
        return await extract_word_text(content_stream, file_name)
    elif "spreadsheetml" in file_type_lower or file_type_lower == "xlsx":
        return await extract_excel_text(content_stream, file_name)
    elif "presentationml" in file_type_lower or file_type_lower == "pptx":
        return await extract_powerpoint_text(content_stream, file_name)
    else:
        return f"[Unsupported Office file type: {file_type}]"
