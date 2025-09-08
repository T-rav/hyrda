"""
Document Processing Service

Handles text extraction from various document formats including:
- PDF files (.pdf)
- Microsoft Office documents (.docx, .xlsx, .pptx)
- Google Workspace documents (handled by Google Drive API export)
- Text files
"""

from io import BytesIO

# Document processing libraries
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word documents
from openpyxl import load_workbook  # openpyxl for Excel
from pptx import Presentation  # python-pptx for PowerPoint


class DocumentProcessor:
    """Service for extracting text content from various document formats."""

    def extract_text(self, content: bytes, mime_type: str) -> str | None:
        """
        Extract text from document bytes based on MIME type.

        Args:
            content: Document content as bytes
            mime_type: MIME type of the document

        Returns:
            Extracted text content, or None if extraction fails
        """
        if mime_type == 'application/pdf':
            return self._extract_pdf_text(content)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self._extract_docx_text(content)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return self._extract_xlsx_text(content)
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return self._extract_pptx_text(content)
        elif mime_type.startswith('text/'):
            return self._extract_text_content(content)
        else:
            print(f"Unsupported MIME type for text extraction: {mime_type}")
            return None

    def _extract_text_content(self, content: bytes) -> str | None:
        """Extract text from plain text content."""
        try:
            return content.decode('utf-8')
        except UnicodeDecodeError:
            return content.decode('latin-1', errors='ignore')

    def _extract_pdf_text(self, pdf_content: bytes) -> str | None:
        """
        Extract text content from PDF bytes using PyMuPDF.

        Args:
            pdf_content: PDF file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PDF from bytes
            pdf_stream = BytesIO(pdf_content)
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

            # Extract text from all pages
            text_content = []
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append(text)

            pdf_document.close()

            # Join all pages with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None

    def _extract_docx_text(self, docx_content: bytes) -> str | None:
        """
        Extract text content from Word document bytes using python-docx.

        Args:
            docx_content: Word document file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Word document from bytes
            docx_stream = BytesIO(docx_content)
            doc = Document(docx_stream)

            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:  # Only add non-empty paragraphs
                    text_content.append(text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_content.append(" | ".join(row_text))

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Word document text: {e}")
            return None

    def _extract_xlsx_text(self, xlsx_content: bytes) -> str | None:
        """
        Extract text content from Excel spreadsheet bytes using openpyxl.

        Args:
            xlsx_content: Excel spreadsheet file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Excel workbook from bytes
            xlsx_stream = BytesIO(xlsx_content)
            workbook = load_workbook(xlsx_stream, read_only=True, data_only=True)

            text_content = []

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]

                # Extract text from all cells with data
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell_value in row:
                        if cell_value is not None:
                            # Convert to string and clean up
                            cell_text = str(cell_value).strip()
                            if cell_text:
                                row_data.append(cell_text)

                    if row_data:  # Only add non-empty rows
                        sheet_content.append(" | ".join(row_data))

                if len(sheet_content) > 1:  # More than just the sheet name
                    text_content.extend(sheet_content)

            workbook.close()

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Excel spreadsheet text: {e}")
            return None

    def _extract_pptx_text(self, pptx_content: bytes) -> str | None:
        """
        Extract text content from PowerPoint presentation bytes using python-pptx.

        Args:
            pptx_content: PowerPoint presentation file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PowerPoint presentation from bytes
            pptx_stream = BytesIO(pptx_content)
            presentation = Presentation(pptx_stream)

            text_content = []

            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]

                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_content.append(shape_text)

                    # Handle tables in slides
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                slide_content.append(" | ".join(row_text))

                if len(slide_content) > 1:  # More than just the slide number
                    text_content.extend(slide_content)

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PowerPoint presentation text: {e}")
            return None
