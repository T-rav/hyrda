#!/usr/bin/env python3
"""
Google Drive Document Ingester - THE ONLY SUPPORTED INGESTION METHOD

This is the sole document ingestion system for the RAG pipeline.
It authenticates with Google Drive, scans folders, and upserts documents
with comprehensive metadata including file paths and permissions.
"""

import argparse
import asyncio
import json
import os
import sys
from datetime import datetime
from io import BytesIO
from pathlib import Path

# Document processing
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word documents

# Google Drive API imports
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from openpyxl import load_workbook  # openpyxl for Excel
from pptx import Presentation  # python-pptx for PowerPoint

# Local imports (assuming we'll use the existing ingestion logic)
sys.path.append(str(Path(__file__).parent.parent / "bot"))

# Load .env file from current directory or parent directories
from dotenv import load_dotenv


def find_and_load_env():
    """Find and load .env file from current directory or parent directories"""
    current_path = Path.cwd()
    for path in [current_path] + list(current_path.parents):
        env_file = path / ".env"
        if env_file.exists():
            load_dotenv(env_file)
            print(f"📄 Loaded environment from: {env_file}")
            return True
    print("⚠️  No .env file found in current directory or parent directories")
    return False

# Load environment variables
find_and_load_env()


from services.embedding_service import chunk_text, create_embedding_provider
from services.vector_service import create_vector_store


class GoogleDriveIngester:
    """Handles Google Drive authentication and document ingestion."""

    # Define the scopes needed for Google Drive API (including shared drives)
    SCOPES = [
        'https://www.googleapis.com/auth/drive.readonly',
        'https://www.googleapis.com/auth/drive.metadata.readonly'
    ]

    def __init__(self, credentials_file: str | None = None, token_file: str | None = None):
        """
        Initialize the Google Drive ingester.

        Args:
            credentials_file: Path to Google OAuth2 credentials JSON file
            token_file: Path to store/retrieve OAuth2 token
        """
        self.credentials_file = credentials_file or 'credentials.json'
        self.token_file = token_file or 'token.json'
        self.service = None
        self.vector_service = None
        self.embedding_service = None

    def authenticate(self) -> bool:
        """
        Authenticate with Google Drive API using OAuth2 or environment variables.

        Returns:
            bool: True if authentication successful, False otherwise
        """
        creds = None

        # Try environment variables first
        client_id = os.getenv('GOOGLE_CLIENT_ID')
        client_secret = os.getenv('GOOGLE_CLIENT_SECRET')
        refresh_token = os.getenv('GOOGLE_REFRESH_TOKEN')

        if client_id and client_secret and refresh_token:
            print("Using credentials from environment variables...")
            try:
                creds = Credentials(
                    token=None,
                    refresh_token=refresh_token,
                    id_token=None,
                    token_uri="https://oauth2.googleapis.com/token",
                    client_id=client_id,
                    client_secret=client_secret,
                    scopes=self.SCOPES
                )
                creds.refresh(Request())
                print("✅ Environment credentials authenticated successfully")
            except Exception as e:
                print(f"❌ Environment credentials failed: {e}")
                return False
        else:
            # Load existing token if available
            if os.path.exists(self.token_file):
                creds = Credentials.from_authorized_user_file(self.token_file, self.SCOPES)

        # If there are no (valid) credentials available, let the user log in
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                try:
                    creds.refresh(Request())
                except Exception as e:
                    print(f"Error refreshing credentials: {e}")
                    return False
            else:
                if not os.path.exists(self.credentials_file):
                    # Create a minimal OAuth2 credentials structure for the flow
                    print("🔐 First time setup - Creating login flow...")

                    # Use a generic OAuth2 client (this would need actual client credentials)
                    # For now, show the user they need to set up OAuth2
                    print("\n❌ No OAuth2 credentials found.")
                    print("📝 To create a login screen like n8n, we need OAuth2 app credentials.")
                    print("\n🚀 One-time setup:")
                    print("1. Go to: https://console.cloud.google.com/")
                    print("2. Create project → Enable Google Drive API → Create OAuth2 Desktop credentials")
                    print("3. Download as 'credentials.json' and put it in this folder")
                    print("4. After that, this will work like n8n - just login once!")
                    return False

                try:
                    flow = InstalledAppFlow.from_client_secrets_file(
                        self.credentials_file, self.SCOPES)
                    creds = flow.run_local_server(port=0)
                except Exception as e:
                    print(f"Error during OAuth flow: {e}")
                    return False

            # Save the credentials for the next run
            try:
                with open(self.token_file, 'w') as token:
                    token.write(creds.to_json())
            except Exception as e:
                print(f"Error saving token: {e}")
                return False

        try:
            self.service = build('drive', 'v3', credentials=creds)
            return True
        except Exception as e:
            print(f"Error building Drive service: {e}")
            return False

    def list_folder_contents(self, folder_id: str, recursive: bool = True, folder_path: str = "") -> list[dict]:
        """
        List all files in a Google Drive folder with comprehensive metadata.

        Args:
            folder_id: Google Drive folder ID
            recursive: Whether to include subfolders
            folder_path: Current folder path for building full paths

        Returns:
            List of file metadata dictionaries with paths and permissions
        """
        if not self.service:
            raise RuntimeError("Not authenticated. Call authenticate() first.")

        files = []

        try:
            # Query for files in the specified folder
            query = f"'{folder_id}' in parents and trashed=false"

            if folder_path == "":  # Only debug for root folder to avoid spam
                print(f"🔍 Querying Google Drive with: {query}")

            # For shared drives, we need different approaches for root vs subfolders
            results = None
            try:
                if folder_path == "":
                    # For root folder, use the broader query approach
                    print("🔄 Using broad shared drive query for root folder...")
                    all_results = self.service.files().list(
                        q="trashed=false",
                        pageSize=1000,
                        includeItemsFromAllDrives=True,
                        supportsAllDrives=True,
                        fields="nextPageToken, files(id, name, mimeType, modifiedTime, size, parents, permissions, owners, createdTime, webViewLink)"
                    ).execute()

                    # Filter to only files that have our folder_id as parent
                    all_files = all_results.get('files', [])
                    filtered_files = [f for f in all_files if folder_id in f.get('parents', [])]
                    results = {'files': filtered_files}
                else:
                    # For subfolders, use the specific parent query with shared drive support
                    print(f"🔄 Using specific parent query for subfolder: {folder_path}")
                    results = self.service.files().list(
                        q=query,
                        pageSize=1000,
                        includeItemsFromAllDrives=True,
                        supportsAllDrives=True,
                        fields="nextPageToken, files(id, name, mimeType, modifiedTime, size, parents, permissions, owners, createdTime, webViewLink)"
                    ).execute()

                # Get files from results (different handling for root vs subfolder)
                if folder_path == "":
                    # Root folder: we already have filtered_files
                    pass
                else:
                    # Subfolder: get files directly from results
                    filtered_files = results.get('files', [])

                # Debug: if no items found for subfolders
                if len(filtered_files) == 0 and folder_path != "":
                    print(f"🔍 DEBUG: No items found for folder_id '{folder_id}' in folder_path '{folder_path}'")
                    print(f"🔍 DEBUG: Query used: {query}")
                    print("🔍 DEBUG: This suggests the folder may be empty or there's a permissions issue")

                if folder_path == "":
                    print(f"✅ Found {len(all_files)} total accessible files, {len(filtered_files)} in target folder")
                    # Debug: let's see a sample of all accessible files to understand the structure
                    print("🔍 DEBUG: Sample of all accessible files (first 10):")
                    for i, f in enumerate(all_files[:10]):
                        parents_str = ', '.join(f.get('parents', []))
                        print(f"   {i+1}. {f.get('name')} (parents: {parents_str})")
                    if len(all_files) > 10:
                        print(f"   ... and {len(all_files) - 10} more files")
                else:
                    print(f"🔍 Subfolder '{folder_path}': Found {len(filtered_files)} items")

            except HttpError as e:
                # Fallback to original query approach
                if folder_path == "":
                    print(f"⚠️  Broad query failed, trying specific parent query: {e}")

                try:
                    results = self.service.files().list(
                        q=query,
                        pageSize=1000,
                        includeItemsFromAllDrives=True,
                        supportsAllDrives=True,
                        fields="nextPageToken, files(id, name, mimeType, modifiedTime, size, parents, permissions, owners, createdTime, webViewLink)"
                    ).execute()
                    if folder_path == "":
                        print(f"✅ Found {len(results.get('files', []))} items using fallback query")
                except HttpError as e2:
                    if folder_path == "":
                        print(f"❌ All query methods failed: {e2}")
                    raise e2

            items = filtered_files

            if folder_path == "" and not items:  # Only debug for root folder
                print(f"📂 Google Drive API returned 0 items for folder {folder_id}")

                # Try to get folder info to see if it exists
                try:
                    # Try regular access first
                    folder_info = self.service.files().get(fileId=folder_id, fields="id,name,mimeType,permissions").execute()
                    print(f"📁 Folder exists: '{folder_info.get('name')}' (Type: {folder_info.get('mimeType')})")
                    permissions = folder_info.get('permissions', [])
                    print(f"🔐 Folder has {len(permissions)} permission entries")
                except HttpError:
                    # Try with shared drive support
                    try:
                        folder_info = self.service.files().get(
                            fileId=folder_id,
                            fields="id,name,mimeType,permissions",
                            supportsAllDrives=True
                        ).execute()
                        print(f"📁 Shared Drive folder exists: '{folder_info.get('name')}' (Type: {folder_info.get('mimeType')})")
                        permissions = folder_info.get('permissions', [])
                        print(f"🔐 Folder has {len(permissions)} permission entries")
                    except HttpError as e2:
                        print(f"❌ Cannot access folder {folder_id}: {e2}")
                        print("💡 This might be a permission issue, invalid folder ID, or shared drive access issue")
                        print("💡 For shared drives, make sure your OAuth app has domain-wide delegation or proper permissions")

            for item in items:
                # Build full path
                current_path = f"{folder_path}/{item['name']}" if folder_path else item['name']

                # Add comprehensive metadata
                item['folder_id'] = folder_id
                item['full_path'] = current_path
                item['folder_path'] = folder_path

                # Get detailed permissions for this file
                try:
                    file_permissions = self.service.files().get(
                        fileId=item['id'],
                        fields="permissions",
                        supportsAllDrives=True
                    ).execute()
                    item['detailed_permissions'] = file_permissions.get('permissions', [])
                except HttpError as perm_error:
                    print(f"Warning: Could not fetch detailed permissions for {item['name']}: {perm_error}")
                    item['detailed_permissions'] = []

                files.append(item)

                # If it's a folder and recursive is enabled, get its contents
                if recursive and item['mimeType'] == 'application/vnd.google-apps.folder':
                    subfolder_files = self.list_folder_contents(
                        item['id'],
                        recursive=True,
                        folder_path=current_path
                    )
                    files.extend(subfolder_files)

        except HttpError as error:
            print(f"An error occurred while listing folder contents: {error}")

        return files

    def download_file_content(self, file_id: str, mime_type: str) -> str | None:
        """
        Download the content of a file from Google Drive.

        Args:
            file_id: Google Drive file ID
            mime_type: MIME type of the file

        Returns:
            File content as string, or None if failed
        """
        if not self.service:
            raise RuntimeError("Not authenticated. Call authenticate() first.")

        try:
            # Handle Google Docs files
            if mime_type == 'application/vnd.google-apps.document':
                request = self.service.files().export_media(
                    fileId=file_id, mimeType='text/plain')
            elif mime_type == 'application/vnd.google-apps.spreadsheet':
                request = self.service.files().export_media(
                    fileId=file_id, mimeType='text/csv')
            elif mime_type == 'application/vnd.google-apps.presentation':
                request = self.service.files().export_media(
                    fileId=file_id, mimeType='text/plain')
            # Handle PDF files with text extraction
            elif mime_type == 'application/pdf':
                request = self.service.files().get_media(fileId=file_id)
                pdf_content = request.execute()
                return self._extract_pdf_text(pdf_content)
            # Handle Office document formats
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':  # .docx
                request = self.service.files().get_media(fileId=file_id)
                docx_content = request.execute()
                return self._extract_docx_text(docx_content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':  # .xlsx
                request = self.service.files().get_media(fileId=file_id)
                xlsx_content = request.execute()
                return self._extract_xlsx_text(xlsx_content)
            elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':  # .pptx
                request = self.service.files().get_media(fileId=file_id)
                pptx_content = request.execute()
                return self._extract_pptx_text(pptx_content)
            # Handle regular text files
            elif mime_type.startswith('text/'):
                request = self.service.files().get_media(fileId=file_id)
            else:
                print(f"Unsupported file type: {mime_type}")
                return None

            # Execute request for Google Apps and text files
            content = request.execute()

            # Decode content based on type
            if isinstance(content, bytes):
                try:
                    return content.decode('utf-8')
                except UnicodeDecodeError:
                    return content.decode('latin-1', errors='ignore')
            else:
                return str(content)

        except HttpError as error:
            print(f"An error occurred downloading file {file_id}: {error}")
            return None

    def _format_permissions(self, permissions: list[dict]) -> dict:
        """
        Format Google Drive permissions into a structured format.

        Args:
            permissions: Raw permissions from Google Drive API

        Returns:
            Formatted permissions dictionary
        """
        formatted = {
            'readers': [],
            'writers': [],
            'owners': [],
            'is_public': False,
            'anyone_can_read': False,
            'anyone_can_write': False
        }

        for perm in permissions:
            role = perm.get('role', 'reader')
            perm_type = perm.get('type', 'user')
            email_address = perm.get('emailAddress', '')
            display_name = perm.get('displayName', email_address)

            # Check for public access
            if perm_type == 'anyone':
                formatted['is_public'] = True
                if role in ['reader', 'commenter']:
                    formatted['anyone_can_read'] = True
                elif role in ['writer', 'editor']:
                    formatted['anyone_can_write'] = True
                    formatted['anyone_can_read'] = True  # Writers can also read

            # Categorize by role
            permission_info = {
                'type': perm_type,
                'email': email_address,
                'display_name': display_name,
                'role': role
            }

            if role == 'owner':
                formatted['owners'].append(permission_info)
            elif role in ['writer', 'editor']:
                formatted['writers'].append(permission_info)
            else:  # reader, commenter
                formatted['readers'].append(permission_info)

        return formatted

    def _get_permissions_summary(self, permissions: list[dict]) -> str:
        """
        Get a simple string summary of Google Drive permissions for Pinecone metadata.

        Args:
            permissions: Raw permissions from Google Drive API

        Returns:
            Simple string summary of permissions
        """
        if not permissions:
            return "no_permissions"

        summary_parts = []
        anyone_access = False
        user_count = 0

        for perm in permissions:
            role = perm.get('role', 'reader')
            perm_type = perm.get('type', 'user')

            if perm_type == 'anyone':
                anyone_access = True
                summary_parts.append(f"anyone_{role}")
            elif perm_type in ['user', 'group']:
                user_count += 1

        if anyone_access and user_count > 0:
            return f"public_plus_{user_count}_users"
        elif anyone_access:
            return "public_access"
        elif user_count > 0:
            return f"private_{user_count}_users"
        else:
            return "restricted"

    def _get_owner_emails(self, owners: list[dict]) -> str:
        """
        Get a simple string of owner emails for Pinecone metadata.

        Args:
            owners: Raw owners from Google Drive API

        Returns:
            Comma-separated string of owner emails
        """
        if not owners:
            return "unknown"

        emails = []
        for owner in owners:
            email = owner.get('emailAddress', '')
            if email:
                emails.append(email)

        return ', '.join(emails) if emails else "unknown"

    def _extract_pdf_text(self, pdf_content: bytes) -> str | None:
        """
        Extract text content from PDF bytes using PyMuPDF.

        Args:
            pdf_content: PDF file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PDF from bytes
            pdf_stream = BytesIO(pdf_content)
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

            # Extract text from all pages
            text_content = []
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append(text)

            pdf_document.close()

            # Join all pages with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None

    def _extract_docx_text(self, docx_content: bytes) -> str | None:
        """
        Extract text content from Word document bytes using python-docx.

        Args:
            docx_content: Word document file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Word document from bytes
            docx_stream = BytesIO(docx_content)
            doc = Document(docx_stream)

            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:  # Only add non-empty paragraphs
                    text_content.append(text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_content.append(" | ".join(row_text))

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Word document text: {e}")
            return None

    def _extract_xlsx_text(self, xlsx_content: bytes) -> str | None:
        """
        Extract text content from Excel spreadsheet bytes using openpyxl.

        Args:
            xlsx_content: Excel spreadsheet file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Excel workbook from bytes
            xlsx_stream = BytesIO(xlsx_content)
            workbook = load_workbook(xlsx_stream, read_only=True, data_only=True)

            text_content = []

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]

                # Extract text from all cells with data
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell_value in row:
                        if cell_value is not None:
                            # Convert to string and clean up
                            cell_text = str(cell_value).strip()
                            if cell_text:
                                row_data.append(cell_text)

                    if row_data:  # Only add non-empty rows
                        sheet_content.append(" | ".join(row_data))

                if len(sheet_content) > 1:  # More than just the sheet name
                    text_content.extend(sheet_content)

            workbook.close()

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Excel spreadsheet text: {e}")
            return None

    def _extract_pptx_text(self, pptx_content: bytes) -> str | None:
        """
        Extract text content from PowerPoint presentation bytes using python-pptx.

        Args:
            pptx_content: PowerPoint presentation file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PowerPoint presentation from bytes
            pptx_stream = BytesIO(pptx_content)
            presentation = Presentation(pptx_stream)

            text_content = []

            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]

                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_content.append(shape_text)

                    # Handle tables in slides
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                slide_content.append(" | ".join(row_text))

                if len(slide_content) > 1:  # More than just the slide number
                    text_content.extend(slide_content)

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PowerPoint presentation text: {e}")
            return None

    async def ingest_files(self, files: list[dict], metadata: dict | None = None) -> tuple[int, int]:
        """
        Ingest files into the vector database.

        Args:
            files: List of file metadata from Google Drive
            metadata: Additional metadata to add to each document

        Returns:
            Tuple of (success_count, error_count)
        """
        # Check that services are properly initialized
        if not self.vector_service:
            raise RuntimeError("Vector service not initialized. Services must be set before calling ingest_files.")
        if not self.embedding_service:
            raise RuntimeError("Embedding service not initialized. Services must be set before calling ingest_files.")

        success_count = 0
        error_count = 0

        for file_info in files:
            try:
                # Skip folders and unsupported file types
                if file_info['mimeType'] == 'application/vnd.google-apps.folder':
                    continue

                print(f"Processing: {file_info['name']}")

                # Download file content
                content = self.download_file_content(file_info['id'], file_info['mimeType'])
                if not content:
                    print(f"Failed to download: {file_info['name']}")
                    error_count += 1
                    continue

                # Prepare comprehensive document metadata
                doc_metadata = {
                    'source': 'google_drive',
                    'file_id': file_info['id'],
                    'file_name': file_info['name'],
                    'full_path': file_info.get('full_path', file_info['name']),
                    'folder_path': file_info.get('folder_path', ''),
                    'mime_type': file_info['mimeType'],
                    'modified_time': file_info.get('modifiedTime'),
                    'created_time': file_info.get('createdTime'),
                    'size': file_info.get('size'),
                    'web_view_link': file_info.get('webViewLink'),
                    'owner_emails': self._get_owner_emails(file_info.get('owners', [])),
                    'permissions_summary': self._get_permissions_summary(file_info.get('detailed_permissions', [])),
                    'ingested_at': datetime.utcnow().isoformat()
                }

                # Add any additional metadata provided
                if metadata:
                    doc_metadata.update(metadata)

                # Chunk the content and generate embeddings
                chunks = chunk_text(content,
                                   chunk_size=self.embedding_service.settings.chunk_size,
                                   chunk_overlap=self.embedding_service.settings.chunk_overlap)

                # Generate embeddings for chunks
                embeddings = await self.embedding_service.get_embeddings(chunks)

                # Prepare metadata for each chunk
                chunk_metadata = []
                for i, chunk in enumerate(chunks):
                    chunk_meta = doc_metadata.copy()
                    chunk_meta['chunk_id'] = f"{file_info['id']}_chunk_{i}"
                    chunk_meta['chunk_index'] = i
                    chunk_meta['total_chunks'] = len(chunks)
                    chunk_metadata.append(chunk_meta)

                # Add documents to vector store
                await self.vector_service.add_documents(
                    texts=chunks,
                    embeddings=embeddings,
                    metadata=chunk_metadata
                )

                print(f" Successfully ingested: {file_info['name']}")
                success_count += 1

            except Exception as e:
                print(f"L Error processing {file_info.get('name', 'unknown')}: {e}")
                error_count += 1

        return success_count, error_count

    async def ingest_folder(self, folder_id: str, recursive: bool = True,
                          metadata: dict | None = None) -> tuple[int, int]:
        """
        Ingest all files from a Google Drive folder.

        Args:
            folder_id: Google Drive folder ID
            recursive: Whether to include subfolders
            metadata: Additional metadata to add to each document

        Returns:
            Tuple of (success_count, error_count)
        """
        print(f"Scanning folder: {folder_id} (recursive: {recursive})")
        files = self.list_folder_contents(folder_id, recursive, folder_path="")

        # Debug output
        if not files:
            print("Found 0 items - folder appears to be empty or inaccessible")
        else:
            folders = [f for f in files if f['mimeType'] == 'application/vnd.google-apps.folder']
            documents = [f for f in files if f['mimeType'] != 'application/vnd.google-apps.folder']

            print(f"Found {len(files)} total items:")
            print(f"  📁 {len(folders)} folders")
            print(f"  📄 {len(documents)} documents")

            if folders:
                print("  Folders found:")
                for folder in folders[:5]:  # Show first 5 folders
                    print(f"    - {folder['full_path']}")
                if len(folders) > 5:
                    print(f"    ... and {len(folders) - 5} more folders")

            if documents:
                print("  Documents found:")
                for doc in documents[:5]:  # Show first 5 documents
                    print(f"    - {doc['full_path']} ({doc['mimeType']})")
                if len(documents) > 5:
                    print(f"    ... and {len(documents) - 5} more documents")

        return await self.ingest_files(files, metadata)


async def main():
    """Main CLI entry point."""
    parser = argparse.ArgumentParser(description="Ingest documents from Google Drive")
    parser.add_argument("--folder-id", required=True, help="Google Drive folder ID")
    parser.add_argument("--credentials", help="Path to Google OAuth2 credentials JSON file")
    parser.add_argument("--token", help="Path to store OAuth2 token")
    parser.add_argument("--recursive", action="store_true", default=True,
                       help="Include subfolders (default: True)")
    parser.add_argument("--metadata", help="Additional metadata as JSON string")
    parser.add_argument("--reauth", action="store_true",
                       help="Force re-authentication (useful for shared drive access)")

    args = parser.parse_args()

    # Parse metadata if provided
    metadata = None
    if args.metadata:
        try:
            metadata = json.loads(args.metadata)
        except json.JSONDecodeError as e:
            print(f"Error parsing metadata JSON: {e}")
            sys.exit(1)

    # Handle re-authentication if requested
    if args.reauth:
        token_file = args.token or 'token.json'
        if os.path.exists(token_file):
            print(f"🔄 Removing existing token file: {token_file}")
            os.remove(token_file)

    # Initialize ingester
    ingester = GoogleDriveIngester(args.credentials, args.token)

    # Authenticate
    print("Authenticating with Google Drive...")
    if not ingester.authenticate():
        print("L Authentication failed")
        sys.exit(1)

    print(" Authentication successful")

    # Initialize services
    print("Initializing vector database and embedding service...")
    try:
        from config.settings import EmbeddingSettings, LLMSettings, VectorSettings

        # Initialize only the settings we need (not full Settings which requires Slack)
        vector_settings = VectorSettings()
        embedding_settings = EmbeddingSettings()
        llm_settings = LLMSettings()

        # Initialize vector store
        vector_store = create_vector_store(vector_settings)
        await vector_store.initialize()
        ingester.vector_service = vector_store

        # Initialize embedding service
        embedding_provider = create_embedding_provider(embedding_settings, llm_settings)
        ingester.embedding_service = embedding_provider

        print("✅ Services initialized successfully")
    except Exception as e:
        print(f"❌ Service initialization failed: {e}")
        sys.exit(1)

    # Ingest folder
    try:
        success_count, error_count = await ingester.ingest_folder(
            args.folder_id,
            recursive=args.recursive,
            metadata=metadata
        )

        print("\n=� Ingestion Summary:")
        print(f" Successfully processed: {success_count}")
        print(f"L Errors: {error_count}")
        print(f"=� Total items: {success_count + error_count}")

    except Exception as e:
        print(f"L Ingestion failed: {e}")
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())
