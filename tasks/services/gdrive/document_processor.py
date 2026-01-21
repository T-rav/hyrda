"""
Document Processing Service

Handles text extraction from various document formats including:
- PDF files (.pdf)
- Microsoft Office documents (.docx, .xlsx, .pptx)
- Google Workspace documents (handled by Google Drive API export)
- Text files
- Audio files (.mp3, .wav, .m4a, .aac, .ogg, .flac, .webm) - transcribed using OpenAI Whisper
- Video files (.mp4, .mov, .avi, .mkv, .webm) - audio extracted and transcribed
"""

import os
import tempfile
from io import BytesIO

# Document processing libraries
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word documents
from openpyxl import load_workbook  # openpyxl for Excel
from pptx import Presentation  # python-pptx for PowerPoint


class DocumentProcessor:
    """Service for extracting text content from various document formats."""

    # Audio MIME types supported for transcription
    AUDIO_MIME_TYPES = {
        "audio/mpeg",  # mp3
        "audio/mp3",  # mp3
        "audio/wav",  # wav
        "audio/x-wav",  # wav
        "audio/wave",  # wav
        "audio/mp4",  # m4a
        "audio/x-m4a",  # m4a
        "audio/aac",  # aac
        "audio/ogg",  # ogg
        "audio/flac",  # flac
        "audio/webm",  # webm
    }

    # Video MIME types supported for transcription
    VIDEO_MIME_TYPES = {
        "video/mp4",  # mp4
        "video/quicktime",  # mov
        "video/x-msvideo",  # avi
        "video/x-matroska",  # mkv
        "video/webm",  # webm
    }

    def __init__(self, openai_api_key: str | None = None):
        """
        Initialize document processor.

        Args:
            openai_api_key: OpenAI API key for transcription. If None, reads from OPENAI_API_KEY env var.
        """
        self.openai_api_key = openai_api_key or os.getenv("OPENAI_API_KEY")

    def extract_text(  # noqa: PLR0911
        self, content: bytes, mime_type: str, filename: str = "file"
    ) -> str | None:
        """
        Extract text from document bytes based on MIME type.

        Args:
            content: Document content as bytes
            mime_type: MIME type of the document
            filename: Original filename (used for audio/video transcription)

        Returns:
            Extracted text content, or None if extraction fails
        """
        if mime_type == "application/pdf":
            return self._extract_pdf_text(content)
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            return self._extract_docx_text(content)
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            return self._extract_xlsx_text(content)
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return self._extract_pptx_text(content)
        elif mime_type.startswith("text/"):
            return self._extract_text_content(content)
        elif mime_type in self.AUDIO_MIME_TYPES:
            return self._transcribe_audio(content, filename)
        elif mime_type in self.VIDEO_MIME_TYPES:
            return self._transcribe_video(content, filename)
        else:
            print(f"Unsupported MIME type for text extraction: {mime_type}")
            return None

    def _extract_text_content(self, content: bytes) -> str | None:
        """Extract text from plain text content."""
        try:
            return content.decode("utf-8")
        except UnicodeDecodeError:
            return content.decode("latin-1", errors="ignore")

    def _extract_pdf_text(self, pdf_content: bytes) -> str | None:
        """
        Extract text content from PDF bytes using PyMuPDF.

        Args:
            pdf_content: PDF file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PDF from bytes
            pdf_stream = BytesIO(pdf_content)
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

            # Extract text from all pages
            text_content = []
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    text_content.append(text)

            pdf_document.close()

            # Join all pages with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PDF text: {e}")
            return None

    def _extract_docx_text(self, docx_content: bytes) -> str | None:
        """
        Extract text content from Word document bytes using python-docx.

        Args:
            docx_content: Word document file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Word document from bytes
            docx_stream = BytesIO(docx_content)
            doc = Document(docx_stream)

            # Extract text from all paragraphs
            text_content = []
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:  # Only add non-empty paragraphs
                    text_content.append(text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_text.append(cell_text)
                    if row_text:
                        text_content.append(" | ".join(row_text))

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Word document text: {e}")
            return None

    def _extract_xlsx_text(self, xlsx_content: bytes) -> str | None:
        """
        Extract text content from Excel spreadsheet bytes using openpyxl.

        Args:
            xlsx_content: Excel spreadsheet file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open Excel workbook from bytes
            xlsx_stream = BytesIO(xlsx_content)
            workbook = load_workbook(xlsx_stream, read_only=True, data_only=True)

            text_content = []

            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]

                # Extract text from all cells with data
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell_value in row:
                        if cell_value is not None:
                            # Convert to string and clean up
                            cell_text = str(cell_value).strip()
                            if cell_text:
                                row_data.append(cell_text)

                    if row_data:  # Only add non-empty rows
                        sheet_content.append(" | ".join(row_data))

                if len(sheet_content) > 1:  # More than just the sheet name
                    text_content.extend(sheet_content)

            workbook.close()

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting Excel spreadsheet text: {e}")
            return None

    def _extract_pptx_text(self, pptx_content: bytes) -> str | None:
        """
        Extract text content from PowerPoint presentation bytes using python-pptx.

        Args:
            pptx_content: PowerPoint presentation file content as bytes

        Returns:
            Extracted text content, or None if extraction fails
        """
        try:
            # Open PowerPoint presentation from bytes
            pptx_stream = BytesIO(pptx_content)
            presentation = Presentation(pptx_stream)

            text_content = []

            # Process each slide
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]

                # Extract text from all shapes on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_text = shape.text.strip()
                        if shape_text:
                            slide_content.append(shape_text)

                    # Handle tables in slides
                    if hasattr(shape, "table"):
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                cell_text = cell.text.strip()
                                if cell_text:
                                    row_text.append(cell_text)
                            if row_text:
                                slide_content.append(" | ".join(row_text))

                if len(slide_content) > 1:  # More than just the slide number
                    text_content.extend(slide_content)

            # Join all content with double newlines
            full_text = "\n\n".join(text_content)
            return full_text if full_text.strip() else None

        except Exception as e:
            print(f"Error extracting PowerPoint presentation text: {e}")
            return None

    def _transcribe_audio(self, audio_content: bytes, filename: str) -> str | None:
        """
        Transcribe audio content to text using OpenAI Whisper API.

        Args:
            audio_content: Audio file content as bytes
            filename: Original filename (used to determine extension)

        Returns:
            Transcribed text content, or None if transcription fails
        """
        if not self.openai_api_key:
            print("OpenAI API key not configured - skipping audio transcription")
            return None

        try:
            from openai import OpenAI

            client = OpenAI(api_key=self.openai_api_key)

            # Determine file extension from filename
            file_ext = filename.split(".")[-1] if "." in filename else "mp3"

            # Create temporary file for audio content
            # (Whisper API requires a file-like object with a name attribute)
            with tempfile.NamedTemporaryFile(
                suffix=f".{file_ext}", delete=False
            ) as temp_file:
                temp_file.write(audio_content)
                temp_file_path = temp_file.name

            try:
                # Transcribe audio using Whisper API
                with open(temp_file_path, "rb") as audio_file:
                    transcript = client.audio.transcriptions.create(
                        model="whisper-1", file=audio_file, response_format="text"
                    )

                # Clean up temporary file
                os.unlink(temp_file_path)

                return transcript if transcript.strip() else None

            except Exception as e:
                # Clean up temporary file on error
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                raise e

        except Exception as e:
            print(f"Error transcribing audio: {e}")
            return None

    def _transcribe_video(self, video_content: bytes, filename: str) -> str | None:
        """
        Transcribe video content to text by extracting audio and using OpenAI Whisper API.

        Args:
            video_content: Video file content as bytes
            filename: Original filename (used to determine extension)

        Returns:
            Transcribed text content, or None if transcription fails
        """
        if not self.openai_api_key:
            print("OpenAI API key not configured - skipping video transcription")
            return None

        try:
            from openai import OpenAI

            client = OpenAI(api_key=self.openai_api_key)

            # Determine file extension from filename
            file_ext = filename.split(".")[-1] if "." in filename else "mp4"

            # Create temporary file for video content
            # Note: For video files, Whisper API will extract audio automatically
            with tempfile.NamedTemporaryFile(
                suffix=f".{file_ext}", delete=False
            ) as temp_file:
                temp_file.write(video_content)
                temp_file_path = temp_file.name

            try:
                # Transcribe video (Whisper API handles audio extraction)
                with open(temp_file_path, "rb") as video_file:
                    transcript = client.audio.transcriptions.create(
                        model="whisper-1", file=video_file, response_format="text"
                    )

                # Clean up temporary file
                os.unlink(temp_file_path)

                return transcript if transcript.strip() else None

            except Exception as e:
                # Clean up temporary file on error
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                raise e

        except Exception as e:
            print(f"Error transcribing video: {e}")
            return None
